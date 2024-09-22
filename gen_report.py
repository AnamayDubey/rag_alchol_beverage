import os
import logging
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def extract_text_from_slide(slide):
    text = ""
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text += shape.text + "\n"
    return text.strip()

def read_description_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read().strip()
        logging.info(f"Successfully read description file: {file_path}")
        return content
    except FileNotFoundError:
        logging.warning(f"Description file not found: {file_path}")
        return "Description file not found."
    except Exception as e:
        logging.error(f"Error reading description file {file_path}: {str(e)}")
        return f"Error reading description file: {str(e)}"

def create_pdf(output_pdf, slide_contents):
    logging.info(f"Creating PDF: {output_pdf}")
    doc = SimpleDocTemplate(output_pdf, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    for i, content in enumerate(slide_contents, 1):
        logging.info(f"Adding content for slide {i} to PDF")
       
        story.append(Paragraph(f"Slide {i}", styles['Heading1']))
        story.append(Spacer(1, 0.2 * inch))

        story.append(Paragraph("Text from PowerPoint:", styles['Heading2']))
        story.append(Paragraph(content['text'], styles['BodyText']))
        story.append(Spacer(1, 0.2 * inch))

        story.append(Paragraph("Image Description:", styles['Heading2']))
        story.append(Paragraph(content['image_description'], styles['BodyText']))
    
        story.append(Spacer(1, inch))

    try:
        doc.build(story)
        logging.info(f"PDF created successfully: {output_pdf}")
    except Exception as e:
        logging.error(f"Error creating PDF: {str(e)}")
        raise

def process_presentation(ppt_file, description_folder, output_pdf):
    logging.info(f"Processing presentation: {ppt_file}")
    try:
        prs = Presentation(ppt_file)
        logging.info(f"Successfully opened presentation: {ppt_file}")
    except Exception as e:
        logging.error(f"Error opening presentation {ppt_file}: {str(e)}")
        raise

    slide_contents = []

    for i, slide in enumerate(prs.slides, 1):
        logging.info(f"Processing slide {i}")
        slide_text = extract_text_from_slide(slide)
        description_file = os.path.join(description_folder, f"slide_{i}.txt")
        image_description = read_description_file(description_file)

        slide_contents.append({
            'text': slide_text,
            'image_description': image_description
        })

    create_pdf(output_pdf, slide_contents)

ppt_file = "alcoholic-beverage-occasions-us-2024__1_ (1).pptx"
description_folder = "descriptions"
output_pdf = "output.pdf"

logging.info("Starting presentation processing")
try:
    process_presentation(ppt_file, description_folder, output_pdf)
    logging.info("Presentation processing completed successfully")
except Exception as e:
    logging.error(f"An error occurred during presentation processing: {str(e)}")